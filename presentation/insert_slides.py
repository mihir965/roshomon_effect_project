"""Insert two appendix slides into the existing pptx:

  1. "How to Read Each Score" — evaluator's plain-language guide
  2. "FPS Weight Justification" — weights + sensitivity analysis

Both slides match the existing deck's visual style (navy `0A2342`,
accent `3182CE`, Montserrat / Inter fonts, 3-column metric table).
The new slides go in after slide 15 (Team Contributions), before the
Additional Results section.

Run:  .venv/bin/python presentation/insert_slides.py
"""

from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
SRC = ROOT / "671 Final Presentation - MVP.pptx"
DST = ROOT / "671 Final Presentation - MVP.pptx"  # overwrite in place

# ── design tokens lifted from slide 5 ────────────────────────────────────────
NAVY = RGBColor(0x0A, 0x23, 0x42)
BLUE = RGBColor(0x31, 0x82, 0xCE)
SLATE = RGBColor(0x2D, 0x37, 0x48)
MUTED = RGBColor(0x71, 0x80, 0x96)
PANEL = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_box(slide, left, top, width, height, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.5)
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, text,
             font="Inter", size=10, bold=False, color=SLATE, align=None,
             fill=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    if fill is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_title_bar(slide, title):
    """Match the existing deck's full-width navy title bar."""
    add_box(slide, 0, 0, Inches(10), Inches(1.25), NAVY)
    add_text(
        slide, Inches(0.62), Inches(0.55), Inches(8.75), Inches(0.46),
        title, font="Montserrat", size=24, bold=True, color=WHITE,
    )


def add_table_header(slide, top, cols):
    """3 (or N) blue header cells. cols = [(left_in, width_in, label), ...]"""
    for left, width, label in cols:
        add_box(slide, Inches(left), top, Inches(width), Inches(0.32), BLUE)
        add_text(
            slide, Inches(left), top, Inches(width), Inches(0.32), label,
            font="Montserrat", size=10, bold=True, color=WHITE,
        )


def add_table_row(slide, top, cells, height_in=0.40, alt=False):
    """cells = [(left_in, width_in, text, bold, color, font), ...]"""
    if alt:
        # alternating row background
        total_left = cells[0][0]
        total_width = cells[-1][0] + cells[-1][1] - total_left
        add_box(slide, Inches(total_left), top, Inches(total_width), Inches(height_in), PANEL)
    for left, width, text, bold, color, font in cells:
        add_text(
            slide, Inches(left), top, Inches(width), Inches(height_in), text,
            font=font, size=9, bold=bold, color=color,
        )


# ── prep ─────────────────────────────────────────────────────────────────────
prs = Presentation(str(SRC))
blank_layout = prs.slide_layouts[10]  # BLANK

# Column geometry mirroring slide 5 exactly
COL1_L, COL1_W = 0.63, 2.42
COL2_L, COL2_W = 3.05, 2.83
COL3_L, COL3_W = 5.88, 3.49


# ─────────────────────────────────────────────────────────────────────────────
# Slide A — "How to Read Each Score"
# ─────────────────────────────────────────────────────────────────────────────
slideA = prs.slides.add_slide(blank_layout)
add_title_bar(slideA, "How to Read Each Score")

# Sub-line under the title bar
add_text(
    slideA, Inches(0.62), Inches(1.05), Inches(8.75), Inches(0.20),
    "Plain-language guide for an evaluator picking a model",
    font="Inter", size=10, color=MUTED,
)

# Table header
header_top = Inches(1.40)
add_table_header(slideA, header_top, [
    (COL1_L, COL1_W, "Metric  —  question it answers"),
    (COL2_L, COL2_W, "What a high value means (≥ 0.7)"),
    (COL3_L, COL3_W, "Most telling for"),
])

rows_a = [
    ('AAS  —  "Did it land on the right answer?"',
     'Same answer as golden, perhaps reworded',
     'Chatbots, search summaries, lookups'),
    ('RAS  —  "Did it think about it the same way?"',
     'Reasoning chain matches in argument structure',
     'High-stakes domains (medical, legal); auditability'),
    ('SLMS  —  "Did it walk through every required step?"',
     'Touched every golden step (asymmetric)',
     'Safety-critical reasoning; no skipped steps'),
    ('CS  —  "Does it tell the same story twice?"',
     'Stable reasoning across repeated runs',
     'Agentic loops; pipelines that re-query'),
    ('DKUS  —  "Does it use the domain vocabulary?"',
     'Uses the required technical terms',
     'Specialist domains (finance, biomedicine, security)'),
    ('FPS  —  "What do I look at first?"',
     'High composite — strong starting recommendation',
     'Leaderboard view; always drill into components'),
]

row_top = Inches(1.72)
for i, (a, b, c) in enumerate(rows_a):
    add_table_row(slideA, row_top, [
        (COL1_L, COL1_W, a, True, NAVY, "Inter SemiBold"),
        (COL2_L, COL2_W, b, False, SLATE, "Inter"),
        (COL3_L, COL3_W, c, False, SLATE, "Inter"),
    ], height_in=0.36, alt=(i % 2 == 1))
    row_top += Inches(0.36)

# Decision shortcut callout
callout_top = Inches(4.10)
add_box(slideA, Inches(0.62), callout_top, Inches(8.75), Inches(1.22), PANEL)
add_text(
    slideA, Inches(0.78), callout_top + Inches(0.10), Inches(8.44), Inches(0.22),
    "Decision shortcut",
    font="Montserrat", size=11.5, bold=True, color=NAVY,
)
shortcuts = [
    "•  High FPS + high RAS + high SLMS  →  recommend with confidence",
    "•  High AAS but low RAS / SLMS  →  classic Rashomon case; right today, brittle tomorrow",
    "•  High everything except CS  →  fine for one-shot use, risky in agentic loops",
    "•  High everything except DKUS  →  fine for general reasoning, suspect in specialist domains",
]
for i, s in enumerate(shortcuts):
    add_text(
        slideA, Inches(0.78), callout_top + Inches(0.34) + Inches(0.20 * i),
        Inches(8.44), Inches(0.20),
        s, font="Inter", size=9.5, color=SLATE,
    )


# ─────────────────────────────────────────────────────────────────────────────
# Slide B — "FPS Weight Justification"
# ─────────────────────────────────────────────────────────────────────────────
slideB = prs.slides.add_slide(blank_layout)
add_title_bar(slideB, "FPS Weight Justification")

add_text(
    slideB, Inches(0.62), Inches(1.05), Inches(8.75), Inches(0.20),
    "Why each metric carries the weight it does — and why the ranking is robust",
    font="Inter", size=10, color=MUTED,
)

# Top table — weight rationale (5 rows)
header_top = Inches(1.40)

rows_b = [
    ('w_RAS = 0.30  (highest)',
     'Full-chain alignment is the most direct test of "did it reason like the golden answer." Captures the Rashomon Effect head-on.'),
    ('w_SLMS = 0.25',
     'Forces the model to walk through each required step, not just produce a chain that aggregates to a similar embedding.'),
    ('w_AAS = 0.20',
     'Answer correctness still matters but every other benchmark already optimises for it. Two equal-AAS models can have very different reasoning quality.'),
    ('w_DKUS = 0.15',
     'Domain vocabulary check. Implementation is substring-based, so synonyms score zero — kept lower until LLM-judge variant lands.'),
    ('w_CS = 0.10  (lowest)',
     'At temperature 0.7 every model scored CS ≥ 0.99 — does not discriminate well in this dataset. Kept as a tie-breaker.'),
]

row_top = Inches(1.72)
COL1B_L, COL1B_W = 0.63, 2.10
COL2B_L, COL2B_W = 2.85, 6.52
add_table_header(slideB, header_top, [
    (COL1B_L, COL1B_W, "Weight"),
    (COL2B_L, COL2B_W, "Justification"),
])

for i, (w, j) in enumerate(rows_b):
    add_table_row(slideB, row_top, [
        (COL1B_L, COL1B_W, w, True, NAVY, "Inter SemiBold"),
        (COL2B_L, COL2B_W, j, False, SLATE, "Inter"),
    ], height_in=0.42, alt=(i % 2 == 1))
    row_top += Inches(0.42)

# Bottom callout — sensitivity
sens_top = row_top + Inches(0.10)
add_box(slideB, Inches(0.62), sens_top, Inches(8.75), Inches(1.30), PANEL)
add_text(
    slideB, Inches(0.78), sens_top + Inches(0.10), Inches(8.44), Inches(0.22),
    "Sensitivity to weight choice  —  rankings are stable",
    font="Montserrat", size=11.5, bold=True, color=NAVY,
)
sens_lines = [
    ("Default / Equal / Reasoning-only schemes",
     "claude-haiku-4-5,  claude-sonnet-4,  claude-opus-4-5    (mistral → 4)"),
    ("AAS-heavy  (w_AAS = 0.50)",
     "claude-sonnet-4,  mistral,  claude-haiku-4-5    (mistral → 2)"),
    ("AAS-only  (w_AAS = 1.00,  the failure mode)",
     "mistral,  claude-sonnet-4,  claude-opus-4-5    (mistral → 1)"),
]
for i, (lhs, rhs) in enumerate(sens_lines):
    y = sens_top + Inches(0.36) + Inches(0.26 * i)
    add_text(slideB, Inches(0.78), y, Inches(3.40), Inches(0.24),
             lhs, font="Inter SemiBold", size=9.5, bold=True, color=NAVY)
    add_text(slideB, Inches(4.20), y, Inches(5.10), Inches(0.24),
             rhs, font="Inter", size=9.5, color=SLATE)


# ─────────────────────────────────────────────────────────────────────────────
# Reorder: move the two new slides (currently at the end) to position 16, 17
# (i.e. right after slide 15 "Team Contributions").
# ─────────────────────────────────────────────────────────────────────────────
xml_slides = prs.slides._sldIdLst
slide_list = list(xml_slides)
# slide_list[-2], slide_list[-1] are our two new slides
new_a, new_b = slide_list[-2], slide_list[-1]
xml_slides.remove(new_a)
xml_slides.remove(new_b)

# Insert at index 15 (which means "after the 15th slide", 0-indexed insertion
# point 15). After re-fetch, list has length 17.
# Note: lxml's insert uses the position in the parent's children list.
parent = xml_slides
# Build current order from the parent again
current = list(parent)
# Insert new_a at index 15, new_b at index 16
parent.insert(15, new_a)
parent.insert(16, new_b)

prs.save(str(DST))
print(f"saved → {DST}")
print("inserted appendix slides at positions 16 and 17")
